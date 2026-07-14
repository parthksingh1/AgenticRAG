"""DOCX, PPTX and XLSX parsers.

Office formats carry their structure explicitly, so these parsers are mostly a
matter of reading it rather than inferring it: Word exposes paragraph styles,
PowerPoint exposes slide and placeholder roles, Excel exposes sheets. That makes
them the highest-fidelity inputs the system accepts, and the parsers preserve
that fidelity rather than flattening everything to text.

Example:
    >>> DocxParser().formats == frozenset({"docx"})
    True
"""

from __future__ import annotations

import io
from typing import Any

from src.core.errors import IngestionError
from src.core.logging import get_logger
from src.ingestion.parsers.base import Parser, register_parser
from src.ingestion.parsers.text import _markdown_table, _title_from_filename
from src.ingestion.types import ParsedDocument, TextBlock
from src.models.document import ChunkKind

log = get_logger(__name__)

#: Cap on rows read per spreadsheet sheet, mirroring the CSV parser.
MAX_SHEET_ROWS = 5_000


class DocxParser(Parser):
    """Reads Word documents, preserving heading levels, tables and lists."""

    name = "python-docx"
    formats = frozenset({"docx"})

    def parse(self, data: bytes, *, filename: str | None = None) -> ParsedDocument:
        """Parse a DOCX into typed blocks."""
        try:
            import docx
        except ImportError as exc:  # pragma: no cover - present in the image
            raise IngestionError("python-docx is not installed") from exc

        try:
            document = docx.Document(io.BytesIO(data))
        except Exception as exc:
            raise IngestionError(f"could not open DOCX: {exc}") from exc

        blocks: list[TextBlock] = []
        section: list[str] = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style = (paragraph.style.name or "").lower()

            if style.startswith("heading"):
                level = _heading_level_from_style(style)
                section = [*section[: level - 1], text]
                blocks.append(
                    TextBlock(
                        text=text,
                        kind=ChunkKind.HEADING,
                        level=level,
                        section_path=tuple(section),
                    )
                )
                continue

            kind = ChunkKind.PROSE
            if "list" in style:
                kind = ChunkKind.LIST
            elif "caption" in style:
                kind = ChunkKind.CAPTION
            elif "quote" in style or "footnote" in style:
                kind = ChunkKind.FOOTNOTE

            blocks.append(TextBlock(text=text, kind=kind, section_path=tuple(section)))

        for table in document.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            markdown = _markdown_table(rows[0], rows[1:]) if rows else ""
            if markdown:
                blocks.append(
                    TextBlock(text=markdown, kind=ChunkKind.TABLE, section_path=tuple(section))
                )

        core = document.core_properties
        return ParsedDocument(
            blocks=tuple(blocks),
            title=(core.title or "").strip() or _title_from_filename(filename),
            parser=self.name,
            metadata={"author": core.author, "created": str(core.created or "")},
        )


class PptxParser(Parser):
    """Reads PowerPoint decks, one section per slide.

    Slide titles become headings so that a chunk always knows which slide it came
    from — the single most useful piece of context when citing a deck, and the
    thing plain text extraction loses first.
    """

    name = "python-pptx"
    formats = frozenset({"pptx"})

    def parse(self, data: bytes, *, filename: str | None = None) -> ParsedDocument:
        """Parse a PPTX into per-slide blocks."""
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover - present in the image
            raise IngestionError("python-pptx is not installed") from exc

        try:
            deck = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise IngestionError(f"could not open PPTX: {exc}") from exc

        blocks: list[TextBlock] = []
        slide_count = 0
        for number, slide in enumerate(deck.slides, start=1):
            slide_count = number
            title = _slide_title(slide) or f"Slide {number}"
            section = (title,)
            blocks.append(
                TextBlock(
                    text=title,
                    kind=ChunkKind.HEADING,
                    level=1,
                    page_number=number,
                    section_path=section,
                )
            )

            for shape in slide.shapes:
                if shape.has_table:
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    markdown = _markdown_table(rows[0], rows[1:]) if rows else ""
                    if markdown:
                        blocks.append(
                            TextBlock(
                                text=markdown,
                                kind=ChunkKind.TABLE,
                                page_number=number,
                                section_path=section,
                            )
                        )
                    continue

                if not getattr(shape, "has_text_frame", False):
                    continue
                body = "\n".join(
                    p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if body and body != title:
                    blocks.append(
                        TextBlock(
                            text=body,
                            kind=ChunkKind.PROSE,
                            page_number=number,
                            section_path=section,
                        )
                    )

            notes = _slide_notes(slide)
            if notes:
                blocks.append(
                    TextBlock(
                        text=notes,
                        kind=ChunkKind.FOOTNOTE,
                        page_number=number,
                        section_path=section,
                        metadata={"speaker_notes": True},
                    )
                )

        return ParsedDocument(
            blocks=tuple(blocks),
            title=_title_from_filename(filename),
            parser=self.name,
            page_count=slide_count,
        )


class XlsxParser(Parser):
    """Reads spreadsheets, one section per sheet."""

    name = "openpyxl"
    formats = frozenset({"xlsx"})

    def parse(self, data: bytes, *, filename: str | None = None) -> ParsedDocument:
        """Parse an XLSX into per-sheet table blocks."""
        try:
            import openpyxl
        except ImportError as exc:  # pragma: no cover - present in the image
            raise IngestionError("openpyxl is not installed") from exc

        try:
            workbook = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        except Exception as exc:
            raise IngestionError(f"could not open XLSX: {exc}") from exc

        blocks: list[TextBlock] = []
        try:
            for sheet in workbook.worksheets:
                rows = _sheet_rows(sheet)
                if not rows:
                    continue
                section = (sheet.title,)
                blocks.append(
                    TextBlock(
                        text=sheet.title,
                        kind=ChunkKind.HEADING,
                        level=1,
                        section_path=section,
                    )
                )
                header, body = rows[0], rows[1:]
                blocks.append(
                    TextBlock(
                        text=f"Sheet '{sheet.title}': columns {', '.join(header)}; "
                        f"{len(body)} rows.",
                        kind=ChunkKind.CAPTION,
                        section_path=section,
                        metadata={"sheet": sheet.title, "columns": header, "row_count": len(body)},
                    )
                )
                for start in range(0, len(body), 50):
                    blocks.append(
                        TextBlock(
                            text=_markdown_table(header, body[start : start + 50]),
                            kind=ChunkKind.TABLE,
                            section_path=section,
                            metadata={"sheet": sheet.title, "row_offset": start},
                        )
                    )
        finally:
            workbook.close()

        return ParsedDocument(
            blocks=tuple(blocks),
            title=_title_from_filename(filename),
            parser=self.name,
            metadata={"sheets": [s.title for s in workbook.worksheets]},
        )


def _sheet_rows(sheet: Any) -> list[list[str]]:
    """Read a sheet's non-empty rows as strings, capped at MAX_SHEET_ROWS."""
    rows: list[list[str]] = []
    for index, raw in enumerate(sheet.iter_rows(values_only=True)):
        if index >= MAX_SHEET_ROWS:
            break
        cells = ["" if value is None else str(value).strip() for value in raw]
        if any(cells):
            rows.append(cells)
    return rows


def _heading_level_from_style(style: str) -> int:
    """Extract the level from a Word heading style name.

    Example:
        >>> _heading_level_from_style("heading 2")
        2
        >>> _heading_level_from_style("heading")
        1
    """
    digits = "".join(ch for ch in style if ch.isdigit())
    if not digits:
        return 1
    return min(max(int(digits), 1), 6)


def _slide_title(slide: Any) -> str | None:
    """The slide's title placeholder text, if it has one."""
    try:
        title = slide.shapes.title
    except AttributeError:
        return None
    if title is None or not getattr(title, "has_text_frame", False):
        return None
    return title.text_frame.text.strip() or None


def _slide_notes(slide: Any) -> str:
    """Speaker notes for a slide, or an empty string."""
    try:
        if not slide.has_notes_slide:
            return ""
        return slide.notes_slide.notes_text_frame.text.strip()
    except (AttributeError, ValueError):
        return ""


docx_parser = register_parser(DocxParser())
pptx_parser = register_parser(PptxParser())
xlsx_parser = register_parser(XlsxParser())
